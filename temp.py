#!/usr/bin/env python3
"""
PPTX Evaluator Test Suite
Automatically tests all .pptx files in test_data_evaluators folder against golden.pptx
"""

import os
import sys
import subprocess
from pathlib import Path
from datetime import datetime
import logging
import xml.etree.ElementTree as ET

# Try importing required libraries
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("ERROR: python-pptx library not found!")
    print("Please install it: pip install python-pptx")
    sys.exit(1)

# ============================================================================
# CONFIGURATION
# ============================================================================
TEST_DATA_DIR = Path("./test_data_evaluators")
GOLDEN_FILE_NAME = "Conservation_Project_Golden_v1.pptx"
TIMESTAMP = datetime.now().strftime("%Y%m%d_%H%M%S")
LOG_FILE = f"pptx_evaluator_test_results_{TIMESTAMP}.log"

# ============================================================================
# COLOR CODES FOR TERMINAL OUTPUT
# ============================================================================
class Colors:
    GREEN = '\033[92m'
    RED = '\033[91m'
    YELLOW = '\033[93m'
    BLUE = '\033[94m'
    MAGENTA = '\033[95m'
    CYAN = '\033[96m'
    BOLD = '\033[1m'
    RESET = '\033[0m'

# ============================================================================
# DEBUG LOGGING SETUP
# ============================================================================
debug_logger = None

def enable_debug_logging():
    """Enable debug logging for detailed comparison output"""
    global debug_logger
    if debug_logger is None:
        debug_logger = logging.getLogger('pptx_debug')
        debug_logger.setLevel(logging.DEBUG)
        
        # Create handler that writes to both file and console
        fh = logging.FileHandler(f'pptx_debug_{TIMESTAMP}.log')
        fh.setLevel(logging.DEBUG)
        
        # Create formatter
        formatter = logging.Formatter('%(message)s')
        fh.setFormatter(formatter)
        
        debug_logger.addHandler(fh)
    return debug_logger

# ============================================================================
# HELPER FUNCTIONS FOR PPTX COMPARISON
# ============================================================================
def get_all_text_shapes(slide):
    """
    Recursively extract all text-containing shapes from a slide,
    including those nested inside GROUP shapes.
    """
    text_shapes = []
    
    def _extract_from_shape(shape):
        # If shape has text_frame, it's a text shape
        if hasattr(shape, 'text_frame'):
            text_shapes.append(shape)
        
        # If it's a GROUP shape, recursively extract from children
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child_shape in shape.shapes:
                _extract_from_shape(child_shape)
    
    # Process all shapes in the slide
    for shape in slide.shapes:
        _extract_from_shape(shape)
    
    return text_shapes

# ============================================================================
# EMBEDDED PPTX EVALUATOR
# ============================================================================
def compare_pptx_files(file1_path, file2_path, **options):
    """
    Compare two PPTX files for equality based on various configurable criteria.
    
    Args:
        file1_path: Path to first PPTX file (test file)
        file2_path: Path to second PPTX file (golden file)
        **options: Dictionary of comparison options
        
    Returns:
        1 if files match according to criteria, 0 otherwise
    """
    prs1 = Presentation(file1_path)
    prs2 = Presentation(file2_path)
    
    # Enable debug logging if requested
    enable_debug = options.get("enable_debug", True)
    if enable_debug:
        enable_debug_logging()
        debug_logger.debug(f"=== COMPARING PPTX FILES ===")
        debug_logger.debug(f"File 1: {file1_path}")
        debug_logger.debug(f"File 2: {file2_path}")
        debug_logger.debug(f"File 1 slides: {len(prs1.slides)}")
        debug_logger.debug(f"File 2 slides: {len(prs2.slides)}")
    
    approximately_tolerance = options.get("approximately_tolerance", 0.005)
    
    def is_approximately_equal(val1, val2, tolerance=approximately_tolerance):
        """
        Compare two values with hybrid tolerance:
        - Percentage tolerance (0.5%) for larger values
        - Absolute tolerance (1000 EMUs ≈ 0.001 inches) for small values
        This handles LibreOffice normalization for both large and tiny shapes.
        """
        # Exact match
        if val1 == val2:
            return True
        
        # Both zero
        if val1 == 0 and val2 == 0:
            return True
        
        # Absolute difference
        abs_diff = abs(val1 - val2)
        
        # ABSOLUTE TOLERANCE: Allow up to 1000 EMUs difference (≈0.001 inches)
        # This handles tiny shapes where percentage tolerance is too strict
        ABSOLUTE_TOLERANCE_EMU = 1000
        if abs_diff <= ABSOLUTE_TOLERANCE_EMU:
            return True
        
        # PERCENTAGE TOLERANCE: For larger shapes, use 0.5% tolerance
        # This handles normal-sized shapes
        if val1 == 0 or val2 == 0:
            return False
        
        percentage_diff = abs_diff / max(abs(val1), abs(val2))
        return percentage_diff <= tolerance
    
    # Extract all examination options
    examine_number_of_slides = options.get("examine_number_of_slides", True)
    examine_shape = options.get("examine_shape", True)
    examine_text = options.get("examine_text", True)
    examine_indent = options.get("examine_indent", True)
    examine_font_name = options.get("examine_font_name", True)
    examine_font_size = options.get("examine_font_size", True)
    examine_font_bold = options.get("examine_font_bold", True)
    examine_font_italic = options.get("examine_font_italic", True)
    examine_color_rgb = options.get("examine_color_rgb", True)
    examine_font_underline = options.get("examine_font_underline", True)
    examine_strike_through = options.get("examine_strike_through", True)
    examine_alignment = options.get("examine_alignment", True)
    examine_title_bottom_position = options.get("examine_title_bottom_position", False)
    examine_table_bottom_position = options.get("examine_table_bottom_position", False)
    examine_right_position = options.get("examine_right_position", False)
    examine_top_position = options.get("examine_top_position", False)
    examine_shape_for_shift_size = options.get("examine_shape_for_shift_size", False)
    examine_image_size = options.get("examine_image_size", False)
    examine_modify_height = options.get("examine_modify_height", False)
    examine_bullets = options.get("examine_bullets", True)
    examine_background_color = options.get("examine_background_color", True)
    examine_note = options.get("examine_note", True)
    
    # Compare the number of slides
    if len(prs1.slides) != len(prs2.slides) and examine_number_of_slides:
        if enable_debug:
            debug_logger.debug(f"MISMATCH: Number of slides differ - File1: {len(prs1.slides)}, File2: {len(prs2.slides)}")
        return 0
    
    slide_idx = 0
    # Compare the content of each slide
    for slide1, slide2 in zip(prs1.slides, prs2.slides):
        slide_idx += 1
        if enable_debug:
            debug_logger.debug(f"--- Comparing Slide {slide_idx} ---")
            debug_logger.debug(f"Slide {slide_idx} - Shapes count: File1={len(slide1.shapes)}, File2={len(slide2.shapes)}")
        
        def get_slide_background_color(slide):
            fill = slide.background.fill
            if fill.type == 1:
                return fill.fore_color.rgb
            elif fill.type == 5:
                master_fill = slide.slide_layout.slide_master.background.fill
                if master_fill.type == 1:
                    return master_fill.fore_color.rgb
                else:
                    return None
            else:
                return None
        
        if get_slide_background_color(slide1) != get_slide_background_color(slide2) and examine_background_color:
            return 0
        
        def get_slide_notes(slide):
            notes_slide = slide.notes_slide
            if notes_slide:
                return notes_slide.notes_text_frame.text
            else:
                return None
        
        if get_slide_notes(slide1).strip() != get_slide_notes(slide2).strip() and examine_note:
            if enable_debug:
                debug_logger.debug(f"    MISMATCH: Slide {slide_idx} - Notes differ:")
                debug_logger.debug(f"      Notes1: '{get_slide_notes(slide1).strip()}'")
                debug_logger.debug(f"      Notes2: '{get_slide_notes(slide2).strip()}'")
            return 0
        
        # Get all text shapes including those inside GROUPs
        text_shapes1 = get_all_text_shapes(slide1)
        text_shapes2 = get_all_text_shapes(slide2)
        
        if enable_debug:
            debug_logger.debug(f"Slide {slide_idx} - Text shapes found: File1={len(text_shapes1)}, File2={len(text_shapes2)}")
        
        # Check if the number of slides is the same
        if len(slide1.shapes) != len(slide2.shapes):
            if enable_debug:
                debug_logger.debug(f"MISMATCH: Slide {slide_idx} - Different number of shapes: File1={len(slide1.shapes)}, File2={len(slide2.shapes)}")
            return 0
        
        # Check if the shapes are the same
        shape_idx = 0
        for shape1, shape2 in zip(slide1.shapes, slide2.shapes):
            shape_idx += 1
            if enable_debug:
                debug_logger.debug(f"  Shape {shape_idx} - Type: {shape1.shape_type} vs {shape2.shape_type}")
                if hasattr(shape1, "text") and hasattr(shape2, "text"):
                    debug_logger.debug(f"  Shape {shape_idx} - Text: '{shape1.text.strip()}' vs '{shape2.text.strip()}'")
                    debug_logger.debug(f"  Shape {shape_idx} - Position: ({shape1.left}, {shape1.top}) vs ({shape2.left}, {shape2.top})")
                    debug_logger.debug(f"  Shape {shape_idx} - Size: ({shape1.width}, {shape1.height}) vs ({shape2.width}, {shape2.height})")
            
            # ============================================================
            # CRITICAL: Shape validation happens BEFORE examine_shape flag checks
            # This ensures explicit validations run even if examine_shape=False
            # ============================================================
            
            # ============== ENHANCEMENT 1: Validate Shape Types ==============
            if shape1.shape_type != shape2.shape_type:
                if enable_debug:
                    debug_logger.debug(f"    MISMATCH: Slide {slide_idx}, Shape {shape_idx} - Shape types differ:")
                    debug_logger.debug(f"      Shape1 type: {shape1.shape_type}")
                    debug_logger.debug(f"      Shape2 type: {shape2.shape_type}")
                return 0
            
            # ============== ENHANCEMENT 2: Validate Picture/Image Shapes ==============
            if shape1.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if enable_debug:
                    debug_logger.debug(f"  Shape {shape_idx} - Validating PICTURE shape")
                
                # Check position with tolerance (allows for LibreOffice normalization)
                if not is_approximately_equal(shape1.left, shape2.left) or \
                   not is_approximately_equal(shape1.top, shape2.top):
                    if enable_debug:
                        debug_logger.debug(f"    MISMATCH: Image position differs beyond tolerance:")
                        debug_logger.debug(f"      Left: {shape1.left} vs {shape2.left} (diff: {abs(shape1.left - shape2.left)})")
                        debug_logger.debug(f"      Top: {shape1.top} vs {shape2.top} (diff: {abs(shape1.top - shape2.top)})")
                    return 0
                
                # Check dimensions with tolerance (allows for LibreOffice normalization)
                if not is_approximately_equal(shape1.width, shape2.width) or \
                   not is_approximately_equal(shape1.height, shape2.height):
                    if enable_debug:
                        debug_logger.debug(f"    MISMATCH: Image dimensions differ beyond tolerance:")
                        debug_logger.debug(f"      Width: {shape1.width} vs {shape2.width} (diff: {abs(shape1.width - shape2.width)})")
                        debug_logger.debug(f"      Height: {shape1.height} vs {shape2.height} (diff: {abs(shape1.height - shape2.height)})")
                    return 0
                
                # Compare image binary content (hash)
                try:
                    import hashlib
                    hash1 = hashlib.md5(shape1.image.blob).hexdigest()
                    hash2 = hashlib.md5(shape2.image.blob).hexdigest()
                    
                    if hash1 != hash2:
                        if enable_debug:
                            debug_logger.debug(f"    MISMATCH: Image content differs")
                            debug_logger.debug(f"      Hash1: {hash1}")
                            debug_logger.debug(f"      Hash2: {hash2}")
                        return 0
                    
                    if enable_debug:
                        debug_logger.debug(f"    ✓ Image validated (hash: {hash1[:8]}...)")
                except Exception as e:
                    if enable_debug:
                        debug_logger.debug(f"    WARNING: Image hash comparison failed - {str(e)}")
            
            
                # ============== ENHANCEMENT: Validate Image Border Properties ==============
                # Check if image has border/line formatting
                if hasattr(shape1, 'line') and hasattr(shape2, 'line'):
                    line1 = shape1.line
                    line2 = shape2.line
                    
                    if enable_debug:
                        debug_logger.debug(f"    Validating image border properties...")
                    
                    # Check border style (dash_style property)
                    # None = no border, 1 = solid/continuous
                    style1 = line1.dash_style
                    style2 = line2.dash_style
                    
                    if style1 != style2:
                        if enable_debug:
                            debug_logger.debug(f"    MISMATCH: Image border style differs:")
                            debug_logger.debug(f"      Style1: {style1} vs Style2: {style2}")
                        return 0
                    
                    # Check border width (in EMUs)
                    # 0.04 inch = 36,576 EMUs (914,400 EMUs per inch)
                    if line1.width != line2.width:
                        # Allow small tolerance for LibreOffice normalization
                        if not is_approximately_equal(line1.width, line2.width, tolerance=0.01):
                            if enable_debug:
                                debug_logger.debug(f"    MISMATCH: Image border width differs:")
                                debug_logger.debug(f"      Width1: {line1.width} EMUs vs Width2: {line2.width} EMUs")
                                if line1.width and line2.width:
                                    debug_logger.debug(f"      Width1: {line1.width/914400:.4f} inches vs Width2: {line2.width/914400:.4f} inches")
                            return 0
                    
                    # Check border color (RGB)
                    # Access color through line.color.rgb if available
                    try:
                        if hasattr(line1.color, 'rgb') and hasattr(line2.color, 'rgb'):
                            color1 = line1.color.rgb
                            color2 = line2.color.rgb
                            
                            if color1 != color2:
                                if enable_debug:
                                    debug_logger.debug(f"    MISMATCH: Image border color differs:")
                                    debug_logger.debug(f"      Color1: {color1} vs Color2: {color2}")
                                return 0
                        elif line1.color.type != line2.color.type:
                            # If color types differ (e.g., RGB vs. Theme), it's a mismatch
                            if enable_debug:
                                debug_logger.debug(f"    MISMATCH: Image border color type differs:")
                                debug_logger.debug(f"      Type1: {line1.color.type} vs Type2: {line2.color.type}")
                            return 0
                    except Exception as e:
                        if enable_debug:
                            debug_logger.debug(f"    WARNING: Border color comparison failed - {str(e)}")
                    
                    if enable_debug:
                        debug_logger.debug(f"    ✓ Image border validated")
            
            # ============== ENHANCEMENT 3: Validate Table Properties ==============
            if shape1.shape_type == MSO_SHAPE_TYPE.TABLE:
                if enable_debug:
                    debug_logger.debug(f"  Shape {shape_idx} - Validating TABLE shape")
                
                # STRICT position check - use absolute tolerance only (max 2000 EMUs ≈ 0.002")
                # Position should be exact or very close, no percentage-based tolerance
                POSITION_TOLERANCE_EMU = 2000  # ≈ 0.002 inches
                
                pos_diff_left = abs(shape1.left - shape2.left)
                pos_diff_top = abs(shape1.top - shape2.top)
                
                if pos_diff_left > POSITION_TOLERANCE_EMU or pos_diff_top > POSITION_TOLERANCE_EMU:
                    if enable_debug:
                        debug_logger.debug(f"    MISMATCH: Table position differs beyond tolerance ({POSITION_TOLERANCE_EMU} EMUs):")
                        debug_logger.debug(f"      Left: {shape1.left} vs {shape2.left} (diff: {pos_diff_left} EMUs)")
                        debug_logger.debug(f"      Top: {shape1.top} vs {shape2.top} (diff: {pos_diff_top} EMUs)")
                        debug_logger.debug(f"      Max allowed: {POSITION_TOLERANCE_EMU} EMUs (≈{POSITION_TOLERANCE_EMU/914400:.4f} inches)")
                    return 0
                
                # Check dimensions with hybrid tolerance (percentage + absolute)
                if not is_approximately_equal(shape1.width, shape2.width) or \
                   not is_approximately_equal(shape1.height, shape2.height):
                    if enable_debug:
                        debug_logger.debug(f"    MISMATCH: Table dimensions differ beyond tolerance:")
                        debug_logger.debug(f"      Width: {shape1.width} vs {shape2.width} (diff: {abs(shape1.width - shape2.width)})")
                        debug_logger.debug(f"      Height: {shape1.height} vs {shape2.height} (diff: {abs(shape1.height - shape2.height)})")
                    return 0
                
                if enable_debug:
                    debug_logger.debug(f"    ✓ Table position and dimensions validated")
                
                if enable_debug:
                    debug_logger.debug(f"    ✓ Table position and dimensions validated")

            
            if examine_title_bottom_position:
                if hasattr(shape1, "text") and hasattr(shape2, "text") and shape1.text == shape2.text:
                    if shape1.text == "Product Comparison" and (shape1.top <= shape2.top or shape1.top < 3600000):
                        return 0
                elif (not is_approximately_equal(shape1.left, shape2.left) or 
                      not is_approximately_equal(shape1.top, shape2.top) or 
                      not is_approximately_equal(shape1.width, shape2.width) or 
                      not is_approximately_equal(shape1.height, shape2.height)):
                    return 0
            
            if examine_table_bottom_position:
                if slide_idx == 3 and shape1.shape_type == 19 and shape2.shape_type == 19:
                    if shape1.top <= shape2.top or shape1.top < 3600000:
                        return 0
                elif (not is_approximately_equal(shape1.left, shape2.left) or 
                      not is_approximately_equal(shape1.top, shape2.top) or 
                      not is_approximately_equal(shape1.width, shape2.width) or 
                      not is_approximately_equal(shape1.height, shape2.height)):
                    return 0
            
            if examine_right_position:
                if slide_idx == 2 and not hasattr(shape1, "text") and not hasattr(shape2, "text"):
                    if shape1.left <= shape2.left or shape1.left < 4320000:
                        return 0
            
            if examine_top_position:
                if slide_idx == 2 and shape1.shape_type == 13 and shape2.shape_type == 13:
                    if shape1.top >= shape2.top or shape1.top > 1980000:
                        return 0
            
            if examine_shape_for_shift_size:
                if (not is_approximately_equal(shape1.left, shape2.left) or 
                    not is_approximately_equal(shape1.top, shape2.top) or 
                    not is_approximately_equal(shape1.width, shape2.width) or 
                    not is_approximately_equal(shape1.height, shape2.height)):
                    if not (hasattr(shape1, "text") and hasattr(shape2, "text") and 
                           shape1.text == shape2.text and 
                           shape1.text == "Elaborate on what you want to discuss."):
                        return 0
            
            # CRITICAL: examine_shape check happens BEFORE examine_modify_height!
            if (not is_approximately_equal(shape1.left, shape2.left) or 
                not is_approximately_equal(shape1.top, shape2.top) or 
                not is_approximately_equal(shape1.width, shape2.width) or 
                not is_approximately_equal(shape1.height, shape2.height)) and examine_shape:
                if enable_debug:
                    debug_logger.debug(f"    MISMATCH: Slide {slide_idx}, Shape {shape_idx} - Shape dimensions differ:")
                    debug_logger.debug(f"      Left: {shape1.left} vs {shape2.left} (equal: {is_approximately_equal(shape1.left, shape2.left)})")
                    debug_logger.debug(f"      Top: {shape1.top} vs {shape2.top} (equal: {is_approximately_equal(shape1.top, shape2.top)})")
                    debug_logger.debug(f"      Width: {shape1.width} vs {shape2.width} (equal: {is_approximately_equal(shape1.width, shape2.width)})")
                    debug_logger.debug(f"      Height: {shape1.height} vs {shape2.height} (equal: {is_approximately_equal(shape1.height, shape2.height)})")
                    if hasattr(shape1, "text") and hasattr(shape2, "text"):
                        debug_logger.debug(f"      Shape text: '{shape1.text.strip()}' vs '{shape2.text.strip()}'")
                return 0
            
            if examine_image_size:
                if shape1.shape_type == 13 and shape2.shape_type == 13:
                    if not is_approximately_equal(shape1.width, shape2.width) or not is_approximately_equal(shape1.height, shape2.height):
                        return 0
                elif (not is_approximately_equal(shape1.left, shape2.left) or 
                      not is_approximately_equal(shape1.top, shape2.top) or 
                      not is_approximately_equal(shape1.width, shape2.width) or 
                      not is_approximately_equal(shape1.height, shape2.height)):
                    return 0
            
            if examine_modify_height:
                if not hasattr(shape1, "text") and not hasattr(shape2, "text") or shape1.shape_type == 5 and shape2.shape_type == 5:
                    if not is_approximately_equal(shape1.height, shape2.height):
                        return 0
                elif (not is_approximately_equal(shape1.left, shape2.left) or 
                      not is_approximately_equal(shape1.top, shape2.top) or 
                      not is_approximately_equal(shape1.width, shape2.width) or 
                      not is_approximately_equal(shape1.height, shape2.height)):
                    return 0
            
            if shape1.shape_type == MSO_SHAPE_TYPE.TABLE:
                table1 = shape1.table
                table2 = shape2.table
                if enable_debug:
                    debug_logger.debug(f"  Shape {shape_idx} - Comparing TABLE with {len(table1.rows)} rows and {len(table1.columns)} columns")
                    debug_logger.debug(f"  Shape {shape_idx} - Table2 has {len(table2.rows)} rows and {len(table2.columns)} columns")
                
                # Check if tables have the same dimensions (SINGLE CHECK - NOT DUPLICATED)
                if len(table1.rows) != len(table2.rows) or len(table1.columns) != len(table2.columns):
                    if enable_debug:
                        debug_logger.debug(f"    MISMATCH: Slide {slide_idx}, Shape {shape_idx} (TABLE) - Table dimensions differ:")
                        debug_logger.debug(f"      Table1: {len(table1.rows)} rows x {len(table1.columns)} columns")
                        debug_logger.debug(f"      Table2: {len(table2.rows)} rows x {len(table2.columns)} columns")
                    return 0
                
                # ============== NEW: Table Cell Text Validation ==============
                if enable_debug:
                    debug_logger.debug(f"    Validating table cell contents...")
                
                def normalize_cell_text(text):
                    """Normalize cell text by removing all whitespace variations"""
                    import re
                    normalized = re.sub(r'\s+', ' ', text)
                    normalized = normalized.strip()
                    return normalized
                
                for row_idx in range(len(table1.rows)):
                    for col_idx in range(len(table1.columns)):
                        cell1 = table1.cell(row_idx, col_idx)
                        cell2 = table2.cell(row_idx, col_idx)
                        
                        text1 = normalize_cell_text(cell1.text)
                        text2 = normalize_cell_text(cell2.text)
                        
                        if text1 != text2:
                            if enable_debug:
                                debug_logger.debug(f"    MISMATCH: Slide {slide_idx}, Shape {shape_idx} (TABLE) - Cell [{row_idx},{col_idx}] text differs:")
                                debug_logger.debug(f"      Cell1: '{text1}'")
                                debug_logger.debug(f"      Cell2: '{text2}'")
                            return 0
                
                if enable_debug:
                    debug_logger.debug(f"    ✓ All table cells validated successfully")
                
                
                
                for row_idx in range(len(table1.rows)):
                    for col_idx in range(len(table1.columns)):
                        cell1 = table1.cell(row_idx, col_idx)
                        cell2 = table2.cell(row_idx, col_idx)
                        
                        # Check if cells have the same number of paragraphs
                        if len(cell1.text_frame.paragraphs) != len(cell2.text_frame.paragraphs):
                            if enable_debug:
                                debug_logger.debug(f"    MISMATCH: Slide {slide_idx}, Shape {shape_idx} (TABLE) - Cell [{row_idx},{col_idx}] - Different number of paragraphs:")
                                debug_logger.debug(f"      Cell1 paragraphs: {len(cell1.text_frame.paragraphs)}")
                                debug_logger.debug(f"      Cell2 paragraphs: {len(cell2.text_frame.paragraphs)}")
                            return 0
                        
                        for para_idx, (para1, para2) in enumerate(zip(cell1.text_frame.paragraphs, cell2.text_frame.paragraphs)):
                            # Check if paragraphs have the same number of runs
                            if len(para1.runs) != len(para2.runs):
                                if enable_debug:
                                    debug_logger.debug(f"    MISMATCH: Slide {slide_idx}, Shape {shape_idx} (TABLE) - Cell [{row_idx},{col_idx}], Para {para_idx} - Different number of runs:")
                                    debug_logger.debug(f"      Para1 runs: {len(para1.runs)}")
                                    debug_logger.debug(f"      Para2 runs: {len(para2.runs)}")
                                return 0
                            
                            for run_idx, (run1, run2) in enumerate(zip(para1.runs, para2.runs)):
                                # Check font color
                                if hasattr(run1.font.color, "rgb") and hasattr(run2.font.color, "rgb"):
                                    if run1.font.color.rgb != run2.font.color.rgb and examine_color_rgb:
                                        if enable_debug:
                                            debug_logger.debug(f"    MISMATCH: Slide {slide_idx}, Shape {shape_idx} (TABLE) - Cell [{row_idx},{col_idx}], Para {para_idx}, Run {run_idx} - Font color differs:")
                                            debug_logger.debug(f"      Color1: {run1.font.color.rgb} vs Color2: {run2.font.color.rgb}")
                                            debug_logger.debug(f"      Cell text: '{cell1.text.strip()}' vs '{cell2.text.strip()}'")
                                            debug_logger.debug(f"      Run text: '{run1.text}' vs '{run2.text}'")
                                        return 0
                                
                                # Check font bold
                                if run1.font.bold != run2.font.bold:
                                    if not ((run1.font.bold is None or run1.font.bold is False) and 
                                           (run2.font.bold is None or run2.font.bold is False)):
                                        if enable_debug:
                                            debug_logger.debug(f"    MISMATCH: Slide {slide_idx}, Shape {shape_idx} (TABLE) - Cell [{row_idx},{col_idx}], Para {para_idx}, Run {run_idx} - Font bold differs:")
                                            debug_logger.debug(f"      Bold1: {run1.font.bold} vs Bold2: {run2.font.bold}")
                                            debug_logger.debug(f"      Run text: '{run1.text}' vs '{run2.text}'")
                                        return 0
                                
                                # Check font italic
                                if run1.font.italic != run2.font.italic:
                                    if not ((run1.font.italic is None or run1.font.italic is False) and 
                                           (run2.font.italic is None or run2.font.italic is False)):
                                        if enable_debug:
                                            debug_logger.debug(f"    MISMATCH: Slide {slide_idx}, Shape {shape_idx} (TABLE) - Cell [{row_idx},{col_idx}], Para {para_idx}, Run {run_idx} - Font italic differs:")
                                            debug_logger.debug(f"      Italic1: {run1.font.italic} vs Italic2: {run2.font.italic}")
                                            debug_logger.debug(f"      Run text: '{run1.text}' vs '{run2.text}'")
                                        return 0
                                
                                # Check font underline
                                if run1.font.underline != run2.font.underline:
                                    if run1.font.underline is not None and run2.font.underline is not None:
                                        if enable_debug:
                                            debug_logger.debug(f"    MISMATCH: Slide {slide_idx}, Shape {shape_idx} (TABLE) - Cell [{row_idx},{col_idx}], Para {para_idx}, Run {run_idx} - Font underline differs:")
                                            debug_logger.debug(f"      Underline1: {run1.font.underline} vs Underline2: {run2.font.underline}")
                                            debug_logger.debug(f"      Run text: '{run1.text}' vs '{run2.text}'")
                                        return 0
                                    if (run1.font.underline is None and run2.font.underline is True) or (run1.font.underline is True and run2.font.underline is None):
                                        if enable_debug:
                                            debug_logger.debug(f"    MISMATCH: Slide {slide_idx}, Shape {shape_idx} (TABLE) - Cell [{row_idx},{col_idx}], Para {para_idx}, Run {run_idx} - Font underline differs (None vs True):")
                                            debug_logger.debug(f"      Underline1: {run1.font.underline} vs Underline2: {run2.font.underline}")
                                            debug_logger.debug(f"      Run text: '{run1.text}' vs '{run2.text}'")
                                        return 0
            
            if hasattr(shape1, "text") and hasattr(shape2, "text"):
                if shape1.text.strip() != shape2.text.strip() and examine_text:
                    return 0
                
                # Check if the number of paragraphs are the same
                if len(shape1.text_frame.paragraphs) != len(shape2.text_frame.paragraphs):
                    if enable_debug:
                        debug_logger.debug(f"    MISMATCH: Slide {slide_idx}, Shape {shape_idx} - Different number of paragraphs:")
                        debug_logger.debug(f"      Shape1 paragraphs: {len(shape1.text_frame.paragraphs)}")
                        debug_logger.debug(f"      Shape2 paragraphs: {len(shape2.text_frame.paragraphs)}")
                    return 0
                
                # Check if the paragraphs are the same
                para_idx = 0
                for para1, para2 in zip(shape1.text_frame.paragraphs, shape2.text_frame.paragraphs):
                    para_idx += 1
                    
                    # Handle alignment comparison
                    if examine_alignment:
                        from pptx.enum.text import PP_ALIGN
                        align1 = para1.alignment
                        align2 = para2.alignment
                        
                        if enable_debug:
                            align1_name = "None" if align1 is None else getattr(align1, 'name', str(align1))
                            align2_name = "None" if align2 is None else getattr(align2, 'name', str(align2))
                            debug_logger.debug(f"    Slide {slide_idx}, Shape {shape_idx}, Para {para_idx} - Alignment: '{align1_name}' vs '{align2_name}'")
                            debug_logger.debug(f"    Slide {slide_idx}, Shape {shape_idx}, Para {para_idx} - Text: '{para1.text}' vs '{para2.text}'")
                        
                        # Convert None to LEFT for comparison
                        if align1 is None:
                            align1 = PP_ALIGN.LEFT
                        if align2 is None:
                            align2 = PP_ALIGN.LEFT
                            
                        if align1 != align2:
                            if enable_debug:
                                align1_final = getattr(align1, 'name', str(align1))
                                align2_final = getattr(align2, 'name', str(align2))
                                debug_logger.debug(f"    MISMATCH: Slide {slide_idx}, Shape {shape_idx}, Para {para_idx} - Alignment differs: '{align1_final}' vs '{align2_final}'")
                            return 0
                    
                    # Check if the runs are the same
                    if para1.text != para2.text and examine_text:
                        return 0
                    if para1.level != para2.level and examine_indent:
                        return 0
                    
                    # Check if the number of runs are the same
                    if len(para1.runs) != len(para2.runs):
                        if enable_debug:
                            debug_logger.debug(f"    MISMATCH: Slide {slide_idx}, Shape {shape_idx}, Para {para_idx} - Different number of runs:")
                            debug_logger.debug(f"      Para1 runs: {len(para1.runs)}")
                            debug_logger.debug(f"      Para2 runs: {len(para2.runs)}")
                        return 0
                    
                    for run1, run2 in zip(para1.runs, para2.runs):
                        # Check if the font properties are the same                        
                        if run1.font.name != run2.font.name and examine_font_name:
                            if enable_debug:
                                debug_logger.debug(f"    MISMATCH: Slide {slide_idx}, Shape {shape_idx}, Para {para_idx} - Font name differs:")
                                debug_logger.debug(f"      Name1: '{run1.font.name}' vs Name2: '{run2.font.name}'")
                                debug_logger.debug(f"      Text: '{run1.text}' vs '{run2.text}'")
                            return 0
                        
                        if run1.font.size != run2.font.size and examine_font_size:
                            if enable_debug:
                                debug_logger.debug(f"    MISMATCH: Slide {slide_idx}, Shape {shape_idx}, Para {para_idx} - Font size differs:")
                                debug_logger.debug(f"      Size1: {run1.font.size} vs Size2: {run2.font.size}")
                                debug_logger.debug(f"      Text: '{run1.text}' vs '{run2.text}'")
                            return 0
                        
                        if run1.font.bold != run2.font.bold and examine_font_bold:
                            # Special handling for None vs False
                            if not ((run1.font.bold is None or run1.font.bold is False) and 
                                   (run2.font.bold is None or run2.font.bold is False)):
                                if enable_debug:
                                    debug_logger.debug(f"    MISMATCH: Slide {slide_idx}, Shape {shape_idx}, Para {para_idx} - Font bold differs:")
                                    debug_logger.debug(f"      Bold1: {run1.font.bold} vs Bold2: {run2.font.bold}")
                                    debug_logger.debug(f"      Text: '{run1.text}' vs '{run2.text}'")
                                return 0
                        
                        if run1.font.italic != run2.font.italic and examine_font_italic:
                            # Special handling for None vs False
                            if not ((run1.font.italic is None or run1.font.italic is False) and 
                                   (run2.font.italic is None or run2.font.italic is False)):
                                if enable_debug:
                                    debug_logger.debug(f"    MISMATCH: Slide {slide_idx}, Shape {shape_idx}, Para {para_idx} - Font italic differs:")
                                    debug_logger.debug(f"      Italic1: {run1.font.italic} vs Italic2: {run2.font.italic}")
                                    debug_logger.debug(f"      Text: '{run1.text}' vs '{run2.text}'")
                                return 0
                        
                        if hasattr(run1.font.color, "rgb") and hasattr(run2.font.color, "rgb"):
                            if run1.font.color.rgb != run2.font.color.rgb and examine_color_rgb:
                                if enable_debug:
                                    debug_logger.debug(f"    MISMATCH: Slide {slide_idx}, Shape {shape_idx}, Para {para_idx} - Font color differs:")
                                    debug_logger.debug(f"      Color1: {run1.font.color.rgb} vs Color2: {run2.font.color.rgb}")
                                    debug_logger.debug(f"      Text: '{run1.text}' vs '{run2.text}'")
                                return 0
                        
                        if run1.font.underline != run2.font.underline and examine_font_underline:
                            if run1.font.underline is not None and run2.font.underline is not None:
                                if enable_debug:
                                    debug_logger.debug(f"    MISMATCH: Slide {slide_idx}, Shape {shape_idx}, Para {para_idx} - Font underline differs:")
                                    debug_logger.debug(f"      Underline1: {run1.font.underline} vs Underline2: {run2.font.underline}")
                                    debug_logger.debug(f"      Text: '{run1.text}' vs '{run2.text}'")
                                return 0
                            if (run1.font.underline is None and run2.font.underline is True) or (run1.font.underline is True and run2.font.underline is None):
                                if enable_debug:
                                    debug_logger.debug(f"    MISMATCH: Slide {slide_idx}, Shape {shape_idx}, Para {para_idx} - Font underline differs (None vs True):")
                                    debug_logger.debug(f"      Underline1: {run1.font.underline} vs Underline2: {run2.font.underline}")
                                    debug_logger.debug(f"      Text: '{run1.text}' vs '{run2.text}'")
                                return 0
                        
                        if run1.font._element.attrib.get('strike', 'noStrike') != run2.font._element.attrib.get('strike', 'noStrike') and examine_strike_through:
                            return 0
                        
                        def _extract_bullets(xml_data):
                            root = ET.fromstring(xml_data)
                            namespaces = {
                                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                            }
                            bullets = []
                            for paragraph in root.findall('.//a:p', namespaces):
                                pPr = paragraph.find('a:pPr', namespaces)
                                if pPr is not None:
                                    lvl = pPr.get('lvl')
                                    buChar = pPr.find('a:buChar', namespaces)
                                    char = buChar.get('char') if buChar is not None else "No Bullet"
                                    buClr = pPr.find('a:buClr/a:srgbClr', namespaces)
                                    color = buClr.get('val') if buClr is not None else "No Color"
                                else:
                                    lvl = "No Level"
                                    char = "No Bullet"
                                    color = "No Color"
                                text = "".join(t.text for t in paragraph.findall('.//a:t', namespaces))
                                
                                # Only add non-empty paragraphs to bullets list
                                if text.strip():
                                    bullets.append((lvl, char, text, color))
                            return bullets
                        
                        def _compare_bullets_with_tolerance(bullets1, bullets2):
                            """Compare bullets with tolerance for minor differences"""
                            if len(bullets1) != len(bullets2):
                                return False
                            
                            for (lvl1, char1, text1, color1), (lvl2, char2, text2, color2) in zip(bullets1, bullets2):
                                # Compare text (most important)
                                if text1 != text2:
                                    return False
                                
                                # Compare bullet character
                                if char1 != char2:
                                    return False
                                
                                # Compare level with tolerance (None and '0' are equivalent)
                                normalized_lvl1 = '0' if lvl1 is None else lvl1
                                normalized_lvl2 = '0' if lvl2 is None else lvl2
                                if normalized_lvl1 != normalized_lvl2:
                                    return False
                            
                            return True
                        
                        if examine_bullets:
                            try:
                                bullets1 = _extract_bullets(run1.part.blob.decode('utf-8'))
                                bullets2 = _extract_bullets(run2.part.blob.decode('utf-8'))
                                
                                # Compare bullets with tolerance for minor differences
                                if not _compare_bullets_with_tolerance(bullets1, bullets2):
                                    return 0
                            except:
                                # If bullet extraction fails, skip bullet comparison
                                pass
        
        # Additional check: compare all text shapes including those in GROUPs
        if examine_alignment and len(text_shapes1) == len(text_shapes2):
            for idx, (tshape1, tshape2) in enumerate(zip(text_shapes1, text_shapes2)):
                if enable_debug:
                    debug_logger.debug(f"  Additional text shape check {idx+1}: '{tshape1.text.strip()[:30]}' vs '{tshape2.text.strip()[:30]}'")
                
                # Compare text content
                if tshape1.text.strip() != tshape2.text.strip() and examine_text:
                    if enable_debug:
                        debug_logger.debug(f"    MISMATCH: Text differs - '{tshape1.text.strip()}' vs '{tshape2.text.strip()}'")
                    return 0
                
                # Check if text shapes have the same number of paragraphs
                if len(tshape1.text_frame.paragraphs) != len(tshape2.text_frame.paragraphs):
                    if enable_debug:
                        debug_logger.debug(f"    MISMATCH: Different number of paragraphs - {len(tshape1.text_frame.paragraphs)} vs {len(tshape2.text_frame.paragraphs)}")
                    return 0
                
                # Compare alignment of each paragraph
                for para_idx, (para1, para2) in enumerate(zip(tshape1.text_frame.paragraphs, tshape2.text_frame.paragraphs)):
                    from pptx.enum.text import PP_ALIGN
                    align1 = para1.alignment
                    align2 = para2.alignment
                    
                    if enable_debug:
                        align1_name = "None" if align1 is None else getattr(align1, 'name', str(align1))
                        align2_name = "None" if align2 is None else getattr(align2, 'name', str(align2))
                        debug_logger.debug(f"    Para {para_idx+1}: Alignment '{align1_name}' vs '{align2_name}'")
                    
                    # Convert None to LEFT for comparison
                    if align1 is None:
                        align1 = PP_ALIGN.LEFT
                    if align2 is None:
                        align2 = PP_ALIGN.LEFT
                        
                    if align1 != align2:
                        if enable_debug:
                            align1_final = getattr(align1, 'name', str(align1))
                            align2_final = getattr(align2, 'name', str(align2))
                            debug_logger.debug(f"    MISMATCH: Alignment differs - '{align1_final}' vs '{align2_final}'")
                        return 0
        elif len(text_shapes1) != len(text_shapes2):
            if enable_debug:
                debug_logger.debug(f"MISMATCH: Different number of text shapes - {len(text_shapes1)} vs {len(text_shapes2)}")
            return 0
    
    if enable_debug:
        debug_logger.debug(f"=== COMPARISON SUCCESSFUL - Files match ===")
    return 1

# ============================================================================
# TEST SUITE FUNCTIONS
# ============================================================================
log_file_handle = None

def log_message(message, color=None):
    """Print message to console and write to log file"""
    global log_file_handle
    
    # Print to console with color
    if color:
        print(f"{color}{message}{Colors.RESET}")
    else:
        print(message)
    
    # Write to log file (without color codes)
    if log_file_handle:
        log_file_handle.write(message + "\n")
        log_file_handle.flush()

def check_prerequisites():
    """Check if all required libraries are installed"""
    log_message("\nChecking prerequisites...", Colors.CYAN)
    
    try:
        from pptx import Presentation
        log_message(f"{Colors.GREEN}✓{Colors.RESET} python-pptx is installed")
        return True
    except ImportError:
        log_message(f"{Colors.RED}✗{Colors.RESET} python-pptx is not installed", Colors.RED)
        log_message("Please install it: pip install python-pptx", Colors.YELLOW)
        return False

def find_test_files():
    """Discover all .pptx files in the test data directory"""
    if not TEST_DATA_DIR.exists():
        log_message(f"{Colors.RED}✗{Colors.RESET} Test data directory not found: {TEST_DATA_DIR}", Colors.RED)
        return None, None
    
    # Find all .pptx files
    all_files = sorted(TEST_DATA_DIR.glob("*.pptx"))
    
    if not all_files:
        log_message(f"{Colors.RED}✗{Colors.RESET} No .pptx files found in {TEST_DATA_DIR}", Colors.RED)
        return None, None
    
    # Separate golden file from test files
    golden_file = None
    test_files = []
    
    for file in all_files:
        if file.name == GOLDEN_FILE_NAME:
            golden_file = file
        else:
            test_files.append(file)
    
    if not golden_file:
        log_message(f"{Colors.RED}✗{Colors.RESET} Golden file '{GOLDEN_FILE_NAME}' not found in {TEST_DATA_DIR}", Colors.RED)
        return None, None
    
    if not test_files:
        log_message(f"{Colors.YELLOW}⚠{Colors.RESET} No test files found (only golden file present)", Colors.YELLOW)
        return golden_file, []
    
    log_message(f"{Colors.GREEN}✓{Colors.RESET} Golden file: {golden_file.name}")
    log_message(f"\nFound {len(test_files)} test file(s):", Colors.CYAN)
    for f in test_files:
        log_message(f"  - {f.name}")
    
    return golden_file, test_files

def run_evaluator(test_file_path, golden_file_path):
    """
    Run the PPTX evaluator on a test file against the golden file.
    
    Returns:
        (result, error_message) where result is 1 for match, 0 for mismatch
    """
    try:
        result = compare_pptx_files(
            str(test_file_path),
            str(golden_file_path),
            enable_debug=True  # Enable debug logging for all comparisons
        )
        return result, None
    except Exception as e:
        return 0, str(e)

def main():
    """Main test suite execution"""
    global log_file_handle
    
    # Print header
    print("=" * 80)
    log_message(f"{Colors.BOLD}PPTX Evaluator Test Suite{Colors.RESET}", Colors.CYAN)
    print("=" * 80)
    print()
    
    # Open log file
    log_file_handle = open(LOG_FILE, 'w', encoding='utf-8')
    
    # Log test start information
    start_time = datetime.now()
    log_message(f"Test started at: {start_time.strftime('%Y-%m-%d %H:%M:%S')}")
    log_message(f"Test data directory: {TEST_DATA_DIR.absolute()}")
    log_message(f"Golden file name: {GOLDEN_FILE_NAME}")
    log_message(f"Log file: {LOG_FILE}")
    log_message(f"Debug log file: pptx_debug_{TIMESTAMP}.log")
    print()
    
    # Check prerequisites
    if not check_prerequisites():
        log_file_handle.close()
        sys.exit(1)
    
    print()
    
    # Find test files
    golden_file, test_files = find_test_files()
    if golden_file is None:
        log_file_handle.close()
        sys.exit(1)
    
    if not test_files:
        log_message("\n" + "=" * 80)
        log_message(f"{Colors.YELLOW}No test files to evaluate{Colors.RESET}", Colors.YELLOW)
        log_message("=" * 80)
        log_file_handle.close()
        return
    
    print()
    print("=" * 80)
    print()
    
    # Run tests
    passed = 0
    failed = 0
    results = []
    
    for idx, test_file in enumerate(test_files, 1):
        log_message(f"[Test {idx}/{len(test_files)}] {test_file.name}", Colors.BOLD)
        log_message(f"File path: {test_file.absolute()}")
        
        result, error = run_evaluator(test_file, golden_file)
        
        if result == 1:
            log_message(f"Result: {Colors.GREEN}✓ PASS{Colors.RESET} - Files match", Colors.GREEN)
            passed += 1
            results.append((test_file.name, "PASS", None))
        else:
            if error:
                log_message(f"Result: {Colors.RED}✗ FAIL{Colors.RESET}", Colors.RED)
                log_message(f"  ERROR: {error}", Colors.RED)
                results.append((test_file.name, "FAIL", error))
            else:
                log_message(f"Result: {Colors.RED}✗ FAIL{Colors.RESET} - Files do not match", Colors.RED)
                log_message(f"  See pptx_debug_{TIMESTAMP}.log for detailed comparison", Colors.YELLOW)
                results.append((test_file.name, "FAIL", "Files do not match (see debug log)"))
            failed += 1
        
        log_message("-" * 80)
        print()
    
    # Print summary
    print("=" * 80)
    log_message(f"{Colors.BOLD}TEST SUMMARY{Colors.RESET}", Colors.CYAN)
    print("=" * 80)
    print()
    
    total_tests = len(test_files)
    pass_rate = (passed / total_tests * 100) if total_tests > 0 else 0
    
    log_message(f"Total Tests: {total_tests}")
    log_message(f"Passed: {Colors.GREEN}{passed}{Colors.RESET}")
    log_message(f"Failed: {Colors.RED}{failed}{Colors.RESET}")
    log_message(f"Pass Rate: {pass_rate:.1f}%")
    print()
    
    # Final verdict
    print("=" * 80)
    if failed == 0:
        log_message(f"{Colors.GREEN}✓ ALL TESTS PASSED{Colors.RESET}", Colors.GREEN)
    else:
        log_message(f"{Colors.RED}⚠ {failed} TEST(S) FAILED - CHECK RESULTS ABOVE{Colors.RESET}", Colors.RED)
    print("=" * 80)
    print()
    
    # Log file location
    end_time = datetime.now()
    duration = (end_time - start_time).total_seconds()
    log_message(f"Test completed at: {end_time.strftime('%Y-%m-%d %H:%M:%S')}")
    log_message(f"Duration: {duration:.2f} seconds")
    log_message(f"\nComplete log saved to: {LOG_FILE}")
    log_message(f"Debug log saved to: pptx_debug_{TIMESTAMP}.log")
    
    # Close log file
    log_file_handle.close()
    
    # Exit with appropriate code
    sys.exit(0 if failed == 0 else 1)

if __name__ == "__main__":
    main()