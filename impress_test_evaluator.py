#!/usr/bin/env python3
"""
PPTX Evaluator Test Suite
Compares .pptx files against a golden reference file.
Uses content-based shape matching to handle internal reordering.
Logs only mismatches/errors for clean output.
"""

import sys
import hashlib
import xml.etree.ElementTree as ET
from pathlib import Path
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("ERROR: python-pptx library not found!")
    print("Please install it: pip install python-pptx")
    sys.exit(1)

# ============================================================================
# CONFIGURATION
# ============================================================================
TEST_DATA_DIR = Path("./test_data_evaluators")
GOLDEN_FILE_NAME = "Conservation_Project_Golden_v1.pptx"
TIMESTAMP = datetime.now().strftime("%Y%m%d_%H%M%S")
LOG_FILE = f"pptx_evaluator_results_{TIMESTAMP}.log"

# ============================================================================
# COLOR CODES FOR TERMINAL OUTPUT
# ============================================================================
class Colors:
    GREEN = '\033[92m'
    RED = '\033[91m'
    YELLOW = '\033[93m'
    CYAN = '\033[96m'
    BOLD = '\033[1m'
    RESET = '\033[0m'

# ============================================================================
# LOGGING
# ============================================================================
log_file_handle = None
mismatch_log = []

def log_mismatch(message):
    """Log a mismatch message to the list for later output"""
    mismatch_log.append(message)

def log_message(message, color=None):
    """Print message to console and write to log file"""
    global log_file_handle
    if color:
        print(f"{color}{message}{Colors.RESET}")
    else:
        print(message)
    if log_file_handle:
        log_file_handle.write(message + "\n")
        log_file_handle.flush()

# ============================================================================
# HELPER FUNCTIONS
# ============================================================================
def get_all_text_shapes(slide):
    """
    Recursively extract all text-containing shapes from a slide,
    including those nested inside GROUP shapes.
    """
    text_shapes = []
    
    def _extract_from_shape(shape):
        if hasattr(shape, 'text_frame'):
            text_shapes.append(shape)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child_shape in shape.shapes:
                _extract_from_shape(child_shape)
    
    for shape in slide.shapes:
        _extract_from_shape(shape)
    
    return text_shapes

def is_approximately_equal(val1, val2, tolerance=0.005):
    """
    Compare two values with hybrid tolerance:
    - Percentage tolerance (0.5%) for larger values
    - Absolute tolerance (1000 EMUs ≈ 0.001 inches) for small values
    """
    if val1 == val2:
        return True
    if val1 == 0 and val2 == 0:
        return True
    
    abs_diff = abs(val1 - val2)
    
    # Absolute tolerance for tiny shapes
    if abs_diff <= 1000:
        return True
    
    # Percentage tolerance for larger shapes
    if val1 == 0 or val2 == 0:
        return False
    
    percentage_diff = abs_diff / max(abs(val1), abs(val2))
    return percentage_diff <= tolerance

def normalize_alignment(alignment):
    """Convert None alignment to LEFT for comparison"""
    return PP_ALIGN.LEFT if alignment is None else alignment

def normalize_cell_text(text):
    """Normalize cell text by removing all whitespace variations"""
    import re
    return re.sub(r'\s+', ' ', text).strip()

def get_slide_background_color(slide):
    """Extract background color from slide"""
    fill = slide.background.fill
    if fill.type == 1:
        return fill.fore_color.rgb
    elif fill.type == 5:
        master_fill = slide.slide_layout.slide_master.background.fill
        if master_fill.type == 1:
            return master_fill.fore_color.rgb
    return None

def get_slide_notes(slide):
    """Extract notes text from slide"""
    notes_slide = slide.notes_slide
    if notes_slide:
        return notes_slide.notes_text_frame.text
    return ""

def extract_bullets(xml_data):
    """Extract bullet information from XML data"""
    root = ET.fromstring(xml_data)
    namespaces = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    }
    bullets = []
    for paragraph in root.findall('.//a:p', namespaces):
        pPr = paragraph.find('a:pPr', namespaces)
        if pPr is not None:
            lvl = pPr.get('lvl')
            buChar = pPr.find('a:buChar', namespaces)
            char = buChar.get('char') if buChar is not None else "No Bullet"
            buClr = pPr.find('a:buClr/a:srgbClr', namespaces)
            color = buClr.get('val') if buClr is not None else "No Color"
        else:
            lvl = "No Level"
            char = "No Bullet"
            color = "No Color"
        text = "".join(t.text for t in paragraph.findall('.//a:t', namespaces))
        if text.strip():
            bullets.append((lvl, char, text, color))
    return bullets

def compare_bullets(bullets1, bullets2):
    """Compare bullets with tolerance for minor differences"""
    if len(bullets1) != len(bullets2):
        return False
    for (lvl1, char1, text1, _), (lvl2, char2, text2, _) in zip(bullets1, bullets2):
        if text1 != text2 or char1 != char2:
            return False
        # Normalize level (None and '0' are equivalent)
        if ('0' if lvl1 is None else lvl1) != ('0' if lvl2 is None else lvl2):
            return False
    return True

def fonts_effectively_equal(val1, val2, treat_none_as_false=True):
    """Compare font boolean properties treating None and False as equivalent"""
    if val1 == val2:
        return True
    if treat_none_as_false:
        return (val1 is None or val1 is False) and (val2 is None or val2 is False)
    return False

# ============================================================================
# SHAPE MATCHING FUNCTIONS
# ============================================================================
def get_shape_signature(shape):
    """
    Generate a signature for a shape to help with matching.
    Returns a tuple of (shape_type, text_content, image_hash)
    """
    shape_type = shape.shape_type
    text_content = ""
    image_hash = ""
    
    # Extract text if available
    if hasattr(shape, 'text'):
        text_content = shape.text.strip()
    
    # Extract image hash if it's a picture
    if shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            image_hash = hashlib.md5(shape.image.blob).hexdigest()
        except Exception:
            pass
    
    # For tables, include cell content as text
    if shape_type == MSO_SHAPE_TYPE.TABLE:
        try:
            table_text = []
            for row_idx in range(len(shape.table.rows)):
                for col_idx in range(len(shape.table.columns)):
                    cell_text = normalize_cell_text(shape.table.cell(row_idx, col_idx).text)
                    if cell_text:
                        table_text.append(cell_text)
            text_content = "|".join(table_text)
        except Exception:
            pass
    
    return (shape_type, text_content, image_hash)


def find_matching_shape(target_shape, candidate_shapes, used_indices):
    """
    Find the best matching shape from candidates for the target shape.
    Returns (matched_shape, matched_index) or (None, -1) if no match found.
    """
    target_sig = get_shape_signature(target_shape)
    target_type, target_text, target_hash = target_sig
    
    # Define text-containing shape types that can be treated as equivalent
    TEXT_SHAPE_TYPES = {
        MSO_SHAPE_TYPE.TEXT_BOX,      # 17
        MSO_SHAPE_TYPE.AUTO_SHAPE,    # 1
        MSO_SHAPE_TYPE.PLACEHOLDER,   # 14
    }
    
    best_match = None
    best_match_idx = -1
    best_score = -1
    
    for idx, candidate in enumerate(candidate_shapes):
        if idx in used_indices:
            continue
        
        cand_sig = get_shape_signature(candidate)
        cand_type, cand_text, cand_hash = cand_sig
        
        score = 0
        
        # Shape type matching logic
        if target_type == cand_type:
            score += 10  # Exact type match bonus
        elif target_type in TEXT_SHAPE_TYPES and cand_type in TEXT_SHAPE_TYPES:
            # Allow text-containing shapes to match each other
            score += 5  # Partial type match bonus
        else:
            # Non-text shapes must match exactly (TABLE, PICTURE, GROUP, etc.)
            continue
        
        # Text content match (high priority)
        if target_text and cand_text:
            if target_text == cand_text:
                score += 100  # Exact text match
            elif target_text in cand_text or cand_text in target_text:
                score += 50  # Partial text match
        elif not target_text and not cand_text:
            score += 20  # Both empty - good for non-text shapes
        
        # Image hash match (for pictures)
        if target_hash and cand_hash:
            if target_hash == cand_hash:
                score += 100  # Exact image match
        
        # Position proximity (lower priority, helps break ties)
        if hasattr(target_shape, 'left') and hasattr(candidate, 'left'):
            pos_diff = abs(target_shape.left - candidate.left) + abs(target_shape.top - candidate.top)
            if pos_diff < 100000:  # Within ~0.1 inch
                score += 10
            elif pos_diff < 500000:  # Within ~0.5 inch
                score += 5
        
        if score > best_score:
            best_score = score
            best_match = candidate
            best_match_idx = idx
    
    return (best_match, best_match_idx)

# ============================================================================
# SHAPE COMPARISON FUNCTIONS
# ============================================================================
def compare_picture_shape(shape1, shape2, ctx):
    """Compare PICTURE shape properties"""
    # Check position
    if not is_approximately_equal(shape1.left, shape2.left) or \
       not is_approximately_equal(shape1.top, shape2.top):
        log_mismatch(f"{ctx} (PICTURE): Position mismatch - test=({shape1.left}, {shape1.top}), golden=({shape2.left}, {shape2.top})")
        return 0
    
    # Check dimensions
    if not is_approximately_equal(shape1.width, shape2.width) or \
       not is_approximately_equal(shape1.height, shape2.height):
        log_mismatch(f"{ctx} (PICTURE): Dimension mismatch - test=({shape1.width}x{shape1.height}), golden=({shape2.width}x{shape2.height})")
        return 0
    
    # Compare image content hash
    try:
        hash1 = hashlib.md5(shape1.image.blob).hexdigest()
        hash2 = hashlib.md5(shape2.image.blob).hexdigest()
        if hash1 != hash2:
            log_mismatch(f"{ctx} (PICTURE): Image content mismatch - hash1={hash1[:8]}, hash2={hash2[:8]}")
            return 0
    except Exception as e:
        log_mismatch(f"{ctx} (PICTURE): Image hash comparison failed - {str(e)}")
    
    # Check border properties
    if hasattr(shape1, 'line') and hasattr(shape2, 'line'):
        line1, line2 = shape1.line, shape2.line
        
        if line1.dash_style != line2.dash_style:
            log_mismatch(f"{ctx} (PICTURE): Border style mismatch - test={line1.dash_style}, golden={line2.dash_style}")
            return 0
        
        if not is_approximately_equal(line1.width or 0, line2.width or 0, tolerance=0.01):
            log_mismatch(f"{ctx} (PICTURE): Border width mismatch - test={line1.width}, golden={line2.width}")
            return 0
        
        try:
            if hasattr(line1.color, 'rgb') and hasattr(line2.color, 'rgb'):
                if line1.color.rgb != line2.color.rgb:
                    log_mismatch(f"{ctx} (PICTURE): Border color mismatch - test={line1.color.rgb}, golden={line2.color.rgb}")
                    return 0
            elif line1.color.type != line2.color.type:
                log_mismatch(f"{ctx} (PICTURE): Border color type mismatch")
                return 0
        except Exception:
            pass
    
    return 1

def compare_table_shape(shape1, shape2, ctx):
    """Compare TABLE shape properties"""
    POSITION_TOLERANCE_EMU = 2000
    
    # Strict position check
    if abs(shape1.left - shape2.left) > POSITION_TOLERANCE_EMU or \
       abs(shape1.top - shape2.top) > POSITION_TOLERANCE_EMU:
        log_mismatch(f"{ctx} (TABLE): Position mismatch - test=({shape1.left}, {shape1.top}), golden=({shape2.left}, {shape2.top})")
        return 0
    
    # Dimension check
    if not is_approximately_equal(shape1.width, shape2.width) or \
       not is_approximately_equal(shape1.height, shape2.height):
        log_mismatch(f"{ctx} (TABLE): Dimension mismatch - test=({shape1.width}x{shape1.height}), golden=({shape2.width}x{shape2.height})")
        return 0
    
    table1, table2 = shape1.table, shape2.table
    
    # Check table dimensions
    if len(table1.rows) != len(table2.rows) or len(table1.columns) != len(table2.columns):
        log_mismatch(f"{ctx} (TABLE): Table size mismatch - test={len(table1.rows)}x{len(table1.columns)}, golden={len(table2.rows)}x{len(table2.columns)}")
        return 0
    
    # Check each cell
    for row_idx in range(len(table1.rows)):
        for col_idx in range(len(table1.columns)):
            cell1 = table1.cell(row_idx, col_idx)
            cell2 = table2.cell(row_idx, col_idx)
            cell_ctx = f"{ctx}, Cell [{row_idx},{col_idx}]"
            
            # Check cell text
            text1 = normalize_cell_text(cell1.text)
            text2 = normalize_cell_text(cell2.text)
            if text1 != text2:
                log_mismatch(f"{cell_ctx}: Text mismatch - test='{text1}', golden='{text2}'")
                return 0
            
            # Check paragraph count
            if len(cell1.text_frame.paragraphs) != len(cell2.text_frame.paragraphs):
                log_mismatch(f"{cell_ctx}: Paragraph count mismatch")
                return 0
            
            # Check each paragraph's runs
            for para_idx, (para1, para2) in enumerate(zip(cell1.text_frame.paragraphs, cell2.text_frame.paragraphs)):
                if len(para1.runs) != len(para2.runs):
                    log_mismatch(f"{cell_ctx}, Para {para_idx}: Run count mismatch")
                    return 0
                
                for run_idx, (run1, run2) in enumerate(zip(para1.runs, para2.runs)):
                    run_ctx = f"{cell_ctx}, Para {para_idx}, Run {run_idx}"
                    
                    # Check font color
                    if hasattr(run1.font.color, "rgb") and hasattr(run2.font.color, "rgb"):
                        if run1.font.color.rgb != run2.font.color.rgb:
                            log_mismatch(f"{run_ctx}: Font color mismatch")
                            return 0
                    
                    # Check font bold
                    if not fonts_effectively_equal(run1.font.bold, run2.font.bold):
                        log_mismatch(f"{run_ctx}: Font bold mismatch - test={run1.font.bold}, golden={run2.font.bold}")
                        return 0
                    
                    # Check font italic
                    if not fonts_effectively_equal(run1.font.italic, run2.font.italic):
                        log_mismatch(f"{run_ctx}: Font italic mismatch")
                        return 0
                    
                    # Check font underline
                    if run1.font.underline != run2.font.underline:
                        if not (run1.font.underline is None and run2.font.underline is None):
                            if (run1.font.underline is None) != (run2.font.underline is None) or \
                               (run1.font.underline is True) != (run2.font.underline is True):
                                log_mismatch(f"{run_ctx}: Font underline mismatch")
                                return 0
    
    return 1

def compare_shape_geometry(shape1, shape2, ctx):
    """Compare shape position and dimensions"""
    if not is_approximately_equal(shape1.left, shape2.left) or \
       not is_approximately_equal(shape1.top, shape2.top) or \
       not is_approximately_equal(shape1.width, shape2.width) or \
       not is_approximately_equal(shape1.height, shape2.height):
        log_mismatch(f"{ctx}: Geometry mismatch - " +
                    f"pos: test=({shape1.left}, {shape1.top}) vs golden=({shape2.left}, {shape2.top}), " +
                    f"size: test=({shape1.width}x{shape1.height}) vs golden=({shape2.width}x{shape2.height})")
        return False
    return True

def compare_text_shape(shape1, shape2, ctx):
    """Compare text shape properties including paragraphs and runs"""
    # Check text content
    if shape1.text.strip() != shape2.text.strip():
        log_mismatch(f"{ctx}: Text content mismatch")
        return 0
    
    # Check paragraph count
    if len(shape1.text_frame.paragraphs) != len(shape2.text_frame.paragraphs):
        log_mismatch(f"{ctx}: Paragraph count mismatch - test={len(shape1.text_frame.paragraphs)}, golden={len(shape2.text_frame.paragraphs)}")
        return 0
    
    # Compare each paragraph
    for para_idx, (para1, para2) in enumerate(zip(shape1.text_frame.paragraphs, shape2.text_frame.paragraphs), 1):
        para_ctx = f"{ctx}, Para {para_idx}"
        
        # Check alignment
        align1 = normalize_alignment(para1.alignment)
        align2 = normalize_alignment(para2.alignment)
        if align1 != align2:
            log_mismatch(f"{para_ctx}: Alignment mismatch - test={align1}, golden={align2}")
            return 0
        
        # Check text
        if para1.text != para2.text:
            log_mismatch(f"{para_ctx}: Text mismatch")
            return 0
        
        # Check indent level
        if para1.level != para2.level:
            log_mismatch(f"{para_ctx}: Indent level mismatch - test={para1.level}, golden={para2.level}")
            return 0
        
        # Check run count
        if len(para1.runs) != len(para2.runs):
            log_mismatch(f"{para_ctx}: Run count mismatch - test={len(para1.runs)}, golden={len(para2.runs)}")
            return 0
        
        # Compare each run
        for run_idx, (run1, run2) in enumerate(zip(para1.runs, para2.runs), 1):
            run_ctx = f"{para_ctx}, Run {run_idx}"
            
            # Font name
            if run1.font.name != run2.font.name:
                log_mismatch(f"{run_ctx}: Font name mismatch - test='{run1.font.name}', golden='{run2.font.name}'")
                return 0
            
            # Font size
            if run1.font.size != run2.font.size:
                log_mismatch(f"{run_ctx}: Font size mismatch - test={run1.font.size}, golden={run2.font.size}")
                return 0
            
            # Font bold
            if not fonts_effectively_equal(run1.font.bold, run2.font.bold):
                log_mismatch(f"{run_ctx}: Font bold mismatch - test={run1.font.bold}, golden={run2.font.bold}")
                return 0
            
            # Font italic
            if not fonts_effectively_equal(run1.font.italic, run2.font.italic):
                log_mismatch(f"{run_ctx}: Font italic mismatch - test={run1.font.italic}, golden={run2.font.italic}")
                return 0
            
            # Font color
            if hasattr(run1.font.color, "rgb") and hasattr(run2.font.color, "rgb"):
                if run1.font.color.rgb != run2.font.color.rgb:
                    log_mismatch(f"{run_ctx}: Font color mismatch - test={run1.font.color.rgb}, golden={run2.font.color.rgb}")
                    return 0
            
            # Font underline
            if run1.font.underline != run2.font.underline:
                if run1.font.underline is not None and run2.font.underline is not None:
                    log_mismatch(f"{run_ctx}: Font underline mismatch")
                    return 0
                if (run1.font.underline is None and run2.font.underline is True) or \
                   (run1.font.underline is True and run2.font.underline is None):
                    log_mismatch(f"{run_ctx}: Font underline mismatch (None vs True)")
                    return 0
            
            # Strikethrough
            strike1 = run1.font._element.attrib.get('strike', 'noStrike')
            strike2 = run2.font._element.attrib.get('strike', 'noStrike')
            if strike1 != strike2:
                log_mismatch(f"{run_ctx}: Strikethrough mismatch - test={strike1}, golden={strike2}")
                return 0
            
            # Bullets
            try:
                bullets1 = extract_bullets(run1.part.blob.decode('utf-8'))
                bullets2 = extract_bullets(run2.part.blob.decode('utf-8'))
                if not compare_bullets(bullets1, bullets2):
                    log_mismatch(f"{run_ctx}: Bullets mismatch")
                    return 0
            except Exception:
                pass  # Skip bullet comparison if extraction fails
    
    return 1

# ============================================================================
# MAIN COMPARISON FUNCTION
# ============================================================================
def compare_pptx_files(file1_path, file2_path):
    """
    Compare two PPTX files for equality using content-based shape matching.
    
    Args:
        file1_path: Path to test file
        file2_path: Path to golden file
        
    Returns:
        1 if files match, 0 otherwise
    """
    global mismatch_log
    mismatch_log = []
    
    prs1 = Presentation(file1_path)
    prs2 = Presentation(file2_path)
    
    # Compare number of slides
    if len(prs1.slides) != len(prs2.slides):
        log_mismatch(f"Slide count mismatch: test={len(prs1.slides)}, golden={len(prs2.slides)}")
        return 0
    
    # Compare each slide
    for slide_idx, (slide1, slide2) in enumerate(zip(prs1.slides, prs2.slides), 1):
        
        # Check background color
        if get_slide_background_color(slide1) != get_slide_background_color(slide2):
            log_mismatch(f"Slide {slide_idx}: Background color mismatch")
            return 0
        
        # Check notes
        notes1 = get_slide_notes(slide1).strip()
        notes2 = get_slide_notes(slide2).strip()
        if notes1 != notes2:
            log_mismatch(f"Slide {slide_idx}: Notes mismatch - test='{notes1[:50]}...', golden='{notes2[:50]}...'")
            return 0
        
        # Check shape count
        if len(slide1.shapes) != len(slide2.shapes):
            log_mismatch(f"Slide {slide_idx}: Shape count mismatch - test={len(slide1.shapes)}, golden={len(slide2.shapes)}")
            return 0
        
        # Match shapes from test file to golden file (content-based matching)
        test_shapes = list(slide1.shapes)
        golden_shapes = list(slide2.shapes)
        used_golden_indices = set()
        
        for test_idx, test_shape in enumerate(test_shapes):
            # Find matching golden shape
            golden_shape, golden_idx = find_matching_shape(test_shape, golden_shapes, used_golden_indices)
            
            if golden_shape is None:
                test_sig = get_shape_signature(test_shape)
                log_mismatch(f"Slide {slide_idx}: No matching shape found for test shape {test_idx + 1} "
                           f"(type={test_sig[0]}, text='{test_sig[1][:30]}...')")
                return 0
            
            used_golden_indices.add(golden_idx)
            ctx = f"Slide {slide_idx}, Shape (test:{test_idx + 1} <-> golden:{golden_idx + 1})"
            
            # Now compare the matched pair
            shape1, shape2 = test_shape, golden_shape
            
            # Validate PICTURE shapes
            if shape1.shape_type == MSO_SHAPE_TYPE.PICTURE:
                result = compare_picture_shape(shape1, shape2, ctx)
                if result == 0:
                    return 0
            
            # Validate TABLE shapes
            if shape1.shape_type == MSO_SHAPE_TYPE.TABLE:
                result = compare_table_shape(shape1, shape2, ctx)
                if result == 0:
                    return 0
            
            # Validate shape dimensions and position
            if not compare_shape_geometry(shape1, shape2, ctx):
                return 0
            
            # Validate text shapes
            if hasattr(shape1, "text") and hasattr(shape2, "text"):
                result = compare_text_shape(shape1, shape2, ctx)
                if result == 0:
                    return 0
        
        # Verify all golden shapes were matched
        if len(used_golden_indices) != len(golden_shapes):
            unmatched = set(range(len(golden_shapes))) - used_golden_indices
            log_mismatch(f"Slide {slide_idx}: {len(unmatched)} golden shape(s) not matched by any test shape")
            return 0
        
        # Additional check: compare all text shapes including those in GROUPs
        text_shapes1 = get_all_text_shapes(slide1)
        text_shapes2 = get_all_text_shapes(slide2)
        
        if len(text_shapes1) != len(text_shapes2):
            log_mismatch(f"Slide {slide_idx}: Text shape count mismatch (including groups) - test={len(text_shapes1)}, golden={len(text_shapes2)}")
            return 0
        
        # Match text shapes by content
        used_golden_text_indices = set()
        
        for test_idx, tshape1 in enumerate(text_shapes1):
            tshape2, golden_text_idx = find_matching_shape(tshape1, text_shapes2, used_golden_text_indices)
            
            if tshape2 is None:
                log_mismatch(f"Slide {slide_idx}: No matching text shape for test TextShape {test_idx + 1} "
                           f"(text='{tshape1.text.strip()[:30]}...')")
                return 0
            
            used_golden_text_indices.add(golden_text_idx)
            ctx = f"Slide {slide_idx}, TextShape (test:{test_idx + 1} <-> golden:{golden_text_idx + 1})"
            
            if tshape1.text.strip() != tshape2.text.strip():
                log_mismatch(f"{ctx}: Text mismatch - test='{tshape1.text.strip()[:30]}...', golden='{tshape2.text.strip()[:30]}...'")
                return 0
            
            if len(tshape1.text_frame.paragraphs) != len(tshape2.text_frame.paragraphs):
                log_mismatch(f"{ctx}: Paragraph count mismatch")
                return 0
            
            for para_idx, (para1, para2) in enumerate(zip(tshape1.text_frame.paragraphs, tshape2.text_frame.paragraphs), 1):
                align1 = normalize_alignment(para1.alignment)
                align2 = normalize_alignment(para2.alignment)
                if align1 != align2:
                    log_mismatch(f"{ctx}, Para {para_idx}: Alignment mismatch - test={align1}, golden={align2}")
                    return 0
    
    return 1

# ============================================================================
# TEST SUITE FUNCTIONS
# ============================================================================
def check_prerequisites():
    """Check if all required libraries are installed"""
    log_message("\nChecking prerequisites...", Colors.CYAN)
    try:
        from pptx import Presentation
        log_message(f"  {Colors.GREEN}✓{Colors.RESET} python-pptx installed")
        return True
    except ImportError:
        log_message(f"  {Colors.RED}✗{Colors.RESET} python-pptx not installed", Colors.RED)
        return False

def find_test_files():
    """Discover all .pptx files in the test data directory"""
    if not TEST_DATA_DIR.exists():
        log_message(f"{Colors.RED}✗ Test data directory not found: {TEST_DATA_DIR}{Colors.RESET}", Colors.RED)
        return None, None
    
    all_files = sorted(TEST_DATA_DIR.glob("*.pptx"))
    if not all_files:
        log_message(f"{Colors.RED}✗ No .pptx files found in {TEST_DATA_DIR}{Colors.RESET}", Colors.RED)
        return None, None
    
    golden_file = None
    test_files = []
    
    for file in all_files:
        if file.name == GOLDEN_FILE_NAME:
            golden_file = file
        else:
            test_files.append(file)
    
    if not golden_file:
        log_message(f"{Colors.RED}✗ Golden file '{GOLDEN_FILE_NAME}' not found{Colors.RESET}", Colors.RED)
        return None, None
    
    log_message(f"  {Colors.GREEN}✓{Colors.RESET} Golden file: {golden_file.name}")
    log_message(f"\nFound {len(test_files)} test file(s):", Colors.CYAN)
    for f in test_files:
        log_message(f"    - {f.name}")
    
    return golden_file, test_files

def run_evaluator(test_file_path, golden_file_path):
    """Run the PPTX evaluator on a test file against the golden file"""
    try:
        result = compare_pptx_files(str(test_file_path), str(golden_file_path))
        return result, mismatch_log[0] if mismatch_log else None
    except Exception as e:
        return 0, str(e)

def main():
    """Main test suite execution"""
    global log_file_handle
    
    print("=" * 70)
    log_message(f"{Colors.BOLD}PPTX Evaluator Test Suite{Colors.RESET}", Colors.CYAN)
    print("=" * 70)
    
    log_file_handle = open(LOG_FILE, 'w', encoding='utf-8')
    
    start_time = datetime.now()
    log_message(f"Started: {start_time.strftime('%Y-%m-%d %H:%M:%S')}")
    log_message(f"Log file: {LOG_FILE}")
    
    if not check_prerequisites():
        log_file_handle.close()
        sys.exit(1)
    
    golden_file, test_files = find_test_files()
    if golden_file is None:
        log_file_handle.close()
        sys.exit(1)
    
    if not test_files:
        log_message(f"\n{Colors.YELLOW}No test files to evaluate{Colors.RESET}", Colors.YELLOW)
        log_file_handle.close()
        return
    
    print("\n" + "-" * 70)
    
    passed = 0
    failed = 0
    
    for idx, test_file in enumerate(test_files, 1):
        log_message(f"\n[{idx}/{len(test_files)}] {test_file.name}", Colors.BOLD)
        
        result, error = run_evaluator(test_file, golden_file)
        
        if result == 1:
            log_message(f"  {Colors.GREEN}✓ PASS{Colors.RESET}", Colors.GREEN)
            passed += 1
        else:
            log_message(f"  {Colors.RED}✗ FAIL{Colors.RESET}", Colors.RED)
            if error:
                log_message(f"    Reason: {error}", Colors.YELLOW)
            failed += 1
    
    # Summary
    print("\n" + "=" * 70)
    log_message(f"{Colors.BOLD}SUMMARY{Colors.RESET}", Colors.CYAN)
    print("=" * 70)
    
    total = len(test_files)
    pass_rate = (passed / total * 100) if total > 0 else 0
    
    log_message(f"Total: {total}  |  Passed: {Colors.GREEN}{passed}{Colors.RESET}  |  Failed: {Colors.RED}{failed}{Colors.RESET}  |  Rate: {pass_rate:.1f}%")
    
    print("=" * 70)
    if failed == 0:
        log_message(f"{Colors.GREEN}✓ ALL TESTS PASSED{Colors.RESET}", Colors.GREEN)
    else:
        log_message(f"{Colors.RED}✗ {failed} TEST(S) FAILED{Colors.RESET}", Colors.RED)
    print("=" * 70)
    
    end_time = datetime.now()
    duration = (end_time - start_time).total_seconds()
    log_message(f"\nCompleted in {duration:.2f}s")
    
    log_file_handle.close()
    sys.exit(0 if failed == 0 else 1)

if __name__ == "__main__":
    main()