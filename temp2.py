#!/usr/bin/env python3
"""
PPTX Evaluator Test Suite - Final Version
Automatically tests all .pptx files in test_data_evaluators folder against golden.pptx
Strictly compares File A vs File B and reports the first specific mismatch found.
"""

import os
import sys
import subprocess
from pathlib import Path
from datetime import datetime
import logging
import xml.etree.ElementTree as ET
import hashlib

# Try importing required libraries
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("ERROR: python-pptx library not found!")
    print("Please install it: pip install python-pptx")
    sys.exit(1)

# ============================================================================
# CONFIGURATION
# ============================================================================
TEST_DATA_DIR = Path("./test_data_evaluators")
GOLDEN_FILE_NAME = "Golder_Solution_marketing_v1.pptx"
TIMESTAMP = datetime.now().strftime("%Y%m%d_%H%M%S")
LOG_FILE = f"pptx_evaluator_test_results_{TIMESTAMP}.log"

# ============================================================================
# COLOR CODES
# ============================================================================
class Colors:
    GREEN = '\033[92m'
    RED = '\033[91m'
    YELLOW = '\033[93m'
    BLUE = '\033[94m'
    BOLD = '\033[1m'
    RESET = '\033[0m'

# ============================================================================
# DEBUG LOGGING
# ============================================================================
debug_logger = None

def get_logger():
    """Get or create the debug logger"""
    global debug_logger
    if debug_logger is None:
        debug_logger = logging.getLogger('pptx_debug')
        debug_logger.setLevel(logging.DEBUG)
        
        # Write to file
        fh = logging.FileHandler(f'pptx_debug_{TIMESTAMP}.log', encoding='utf-8')
        fh.setLevel(logging.DEBUG)
        formatter = logging.Formatter('%(message)s')
        fh.setFormatter(formatter)
        debug_logger.addHandler(fh)
    return debug_logger

# ============================================================================
# HELPER FUNCTIONS
# ============================================================================
def is_approximately_equal(val1, val2, tolerance=0.005):
    """Hybrid tolerance comparison for coordinates/sizes."""
    v1 = val1 if val1 is not None else 0
    v2 = val2 if val2 is not None else 0

    if v1 == v2: return True
    
    abs_diff = abs(v1 - v2)
    # 1000 EMUs is approx 0.001 inches (handling tiny rounding errors)
    if abs_diff <= 1000:
        return True
        
    if v1 == 0 or v2 == 0: return False
    percentage_diff = abs_diff / max(abs(v1), abs(v2))
    return percentage_diff <= tolerance

# ============================================================================
# CORE COMPARISON LOGIC
# ============================================================================
def compare_pptx_files(file1_path, file2_path):
    """
    Compare two PPTX files strictly.
    Returns: (bool, str) -> (Passed?, Failure Reason)
    """
    logger = get_logger()
    logger.debug(f"\n=== COMPARING: {os.path.basename(file1_path)} vs {os.path.basename(file2_path)} ===")
    
    try:
        prs1 = Presentation(file1_path)
        prs2 = Presentation(file2_path)
    except Exception as e:
        msg = f"CRITICAL ERROR loading files: {e}"
        logger.debug(msg)
        return False, msg

    # 1. Check Slide Count
    if len(prs1.slides) != len(prs2.slides):
        msg = f"Slide count differs: Test={len(prs1.slides)}, Golden={len(prs2.slides)}"
        logger.debug(f"MISMATCH: {msg}")
        return False, msg

    # Iterate Slides
    for slide_idx, (slide1, slide2) in enumerate(zip(prs1.slides, prs2.slides), 1):
        logger.debug(f"--- Checking Slide {slide_idx} ---")

        # 2. Check Background Color
        def get_bg_color(slide):
            fill = slide.background.fill
            try:
                if fill.type == 1: return fill.fore_color.rgb
                elif fill.type == 5: # Master inheritance
                    return slide.slide_layout.slide_master.background.fill.fore_color.rgb
            except: pass
            return None

        bg1, bg2 = get_bg_color(slide1), get_bg_color(slide2)
        if bg1 != bg2:
            msg = f"Slide {slide_idx}: Background color differs ({bg1} vs {bg2})"
            logger.debug(f"MISMATCH: {msg}")
            return False, msg

        # 3. Check Notes
        notes1 = slide1.notes_slide.notes_text_frame.text if slide1.notes_slide else ""
        notes2 = slide2.notes_slide.notes_text_frame.text if slide2.notes_slide else ""
        if notes1.strip() != notes2.strip():
            msg = f"Slide {slide_idx}: Notes text differs"
            logger.debug(f"MISMATCH: {msg}")
            logger.debug(f"   Note1: {notes1.strip()[:20]}...")
            logger.debug(f"   Note2: {notes2.strip()[:20]}...")
            return False, msg

        # 4. Check Shape Count
        if len(slide1.shapes) != len(slide2.shapes):
            msg = f"Slide {slide_idx}: Shape count differs (Test={len(slide1.shapes)}, Golden={len(slide2.shapes)})"
            logger.debug(f"MISMATCH: {msg}")
            return False, msg

        # Iterate Shapes
        for shape_idx, (shape1, shape2) in enumerate(zip(slide1.shapes, slide2.shapes), 1):
            
            # 5. Check Shape Type
            if shape1.shape_type != shape2.shape_type:
                msg = f"Slide {slide_idx}, Shape {shape_idx}: Type differs ({shape1.shape_type} vs {shape2.shape_type})"
                logger.debug(f"MISMATCH: {msg}")
                return False, msg

            # 6. Check Position & Dimensions (General check for ALL shapes)
            is_table = shape1.shape_type == MSO_SHAPE_TYPE.TABLE
            pos_tolerance = 2000 if is_table else 1000 # EMUs
            
            # Position check
            if abs(shape1.left - shape2.left) > pos_tolerance or abs(shape1.top - shape2.top) > pos_tolerance:
                 # Fallback to percentage for non-tables if absolute failed
                if is_table or (not is_approximately_equal(shape1.left, shape2.left) or not is_approximately_equal(shape1.top, shape2.top)):
                    msg = f"Slide {slide_idx}, Shape {shape_idx}: Position differs"
                    logger.debug(f"MISMATCH: {msg}")
                    logger.debug(f"   Left: {shape1.left} vs {shape2.left}, Top: {shape1.top} vs {shape2.top}")
                    return False, msg

            # Size check
            if not is_approximately_equal(shape1.width, shape2.width) or not is_approximately_equal(shape1.height, shape2.height):
                msg = f"Slide {slide_idx}, Shape {shape_idx}: Dimensions differ"
                logger.debug(f"MISMATCH: {msg}")
                logger.debug(f"   Width: {shape1.width} vs {shape2.width}, Height: {shape1.height} vs {shape2.height}")
                return False, msg

            # 7. Check Images (Content and Borders)
            if shape1.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Content Hash
                try:
                    hash1 = hashlib.md5(shape1.image.blob).hexdigest()
                    hash2 = hashlib.md5(shape2.image.blob).hexdigest()
                    if hash1 != hash2:
                        msg = f"Slide {slide_idx}, Shape {shape_idx}: Image content (binary) differs"
                        logger.debug(f"MISMATCH: {msg}")
                        return False, msg
                except Exception as e:
                    logger.debug(f"WARNING: Image hash check failed: {e}")

                # Image Borders
                if hasattr(shape1, 'line') and hasattr(shape2, 'line'):
                    l1, l2 = shape1.line, shape2.line
                    
                    if l1.dash_style != l2.dash_style:
                        msg = f"Slide {slide_idx}, Shape {shape_idx}: Image border style differs"
                        logger.debug(f"MISMATCH: {msg}")
                        return False, msg
                    
                    w1 = l1.width if l1.width else 0
                    w2 = l2.width if l2.width else 0
                    if abs(w1 - w2) > 100: 
                        msg = f"Slide {slide_idx}, Shape {shape_idx}: Image border width differs"
                        logger.debug(f"MISMATCH: {msg}")
                        return False, msg

                    try:
                        # RGB Check
                        if hasattr(l1.color, 'rgb') and hasattr(l2.color, 'rgb'):
                            if l1.color.rgb != l2.color.rgb:
                                msg = f"Slide {slide_idx}, Shape {shape_idx}: Image border color differs ({l1.color.rgb} vs {l2.color.rgb})"
                                logger.debug(f"MISMATCH: {msg}")
                                return False, msg
                        # Type Check
                        elif l1.color.type != l2.color.type:
                            msg = f"Slide {slide_idx}, Shape {shape_idx}: Image border color type differs"
                            logger.debug(f"MISMATCH: {msg}")
                            return False, msg
                    except: pass

            # 8. Check Tables
            if shape1.shape_type == MSO_SHAPE_TYPE.TABLE:
                tbl1, tbl2 = shape1.table, shape2.table
                if len(tbl1.rows) != len(tbl2.rows) or len(tbl1.columns) != len(tbl2.columns):
                    msg = f"Slide {slide_idx}, Shape {shape_idx}: Table grid differs ({len(tbl1.rows)}x{len(tbl1.columns)} vs {len(tbl2.rows)}x{len(tbl2.columns)})"
                    logger.debug(f"MISMATCH: {msg}")
                    return False, msg
                
                # Check all cells
                for r in range(len(tbl1.rows)):
                    for c in range(len(tbl1.columns)):
                        c1, c2 = tbl1.cell(r, c), tbl2.cell(r, c)
                        passed, reason = compare_text_frames(c1.text_frame, c2.text_frame, logger, f"Table Cell [{r},{c}]")
                        if not passed:
                            msg = f"Slide {slide_idx}, Table Shape {shape_idx}: {reason}"
                            logger.debug(f"MISMATCH: {msg}")
                            return False, msg

            # 9. Check Text Boxes (Content and Formatting)
            if hasattr(shape1, "text_frame") and hasattr(shape2, "text_frame"):
                passed, reason = compare_text_frames(shape1.text_frame, shape2.text_frame, logger, f"Shape {shape_idx}")
                if not passed:
                    msg = f"Slide {slide_idx}, {reason}"
                    logger.debug(f"MISMATCH: {msg}")
                    return False, msg

    logger.debug("=== SUCCESS: Files Match ===")
    return True, "Files Match"

def compare_text_frames(tf1, tf2, logger, context_str):
    """Deep comparison of two text frames. Returns (bool, reason_str)"""
    
    # Simple full text check first
    if tf1.text.strip() != tf2.text.strip():
        return False, f"{context_str} text content differs ('{tf1.text.strip()[:15]}...' vs '{tf2.text.strip()[:15]}...')"

    if len(tf1.paragraphs) != len(tf2.paragraphs):
        return False, f"{context_str} paragraph count differs"

    for p_idx, (p1, p2) in enumerate(zip(tf1.paragraphs, tf2.paragraphs)):
        # Alignment
        a1 = p1.alignment if p1.alignment else PP_ALIGN.LEFT
        a2 = p2.alignment if p2.alignment else PP_ALIGN.LEFT
        if a1 != a2:
            return False, f"{context_str} Para {p_idx} alignment differs ({a1} vs {a2})"
            
        # Bullet Points / Level
        if p1.level != p2.level:
            return False, f"{context_str} Para {p_idx} indent level differs ({p1.level} vs {p2.level})"

        # Runs (Formatting)
        if len(p1.runs) != len(p2.runs):
            return False, f"{context_str} Para {p_idx} run count differs (formatting split mismatch)"

        for r_idx, (r1, r2) in enumerate(zip(p1.runs, p2.runs)):
            # Font Name
            if r1.font.name != r2.font.name:
                return False, f"{context_str} Run {r_idx} font name differs ({r1.font.name} vs {r2.font.name})"
            
            # Font Size
            if r1.font.size != r2.font.size:
                return False, f"{context_str} Run {r_idx} font size differs ({r1.font.size} vs {r2.font.size})"
                
            # Bold/Italic/Underline
            def bool_prop_match(b1, b2):
                return bool(b1) == bool(b2)

            if not bool_prop_match(r1.font.bold, r2.font.bold): 
                return False, f"{context_str} Run {r_idx} bold setting differs"
            if not bool_prop_match(r1.font.italic, r2.font.italic): 
                return False, f"{context_str} Run {r_idx} italic setting differs"
            if not bool_prop_match(r1.font.underline, r2.font.underline): 
                return False, f"{context_str} Run {r_idx} underline setting differs"

            # Font Color
            try:
                if hasattr(r1.font.color, 'rgb') and hasattr(r2.font.color, 'rgb'):
                    if r1.font.color.rgb != r2.font.color.rgb:
                        return False, f"{context_str} Run {r_idx} font color differs"
            except: pass 

    return True, "Match"

# ============================================================================
# MAIN EXECUTION
# ============================================================================
def main():
    log_file_handle = open(LOG_FILE, 'w', encoding='utf-8')
    
    def log(msg, color=None):
        if color: print(f"{color}{msg}{Colors.RESET}")
        else: print(msg)
        log_file_handle.write(msg + "\n")

    log(f"Starting PPTX Evaluation at {datetime.now()}", Colors.GREEN)
    log(f"Data Dir: {TEST_DATA_DIR}")
    
    # Locate files
    if not TEST_DATA_DIR.exists():
        log("Test directory not found!", Colors.RED)
        sys.exit(1)
        
    all_files = sorted(TEST_DATA_DIR.glob("*.pptx"))
    golden_file = next((f for f in all_files if f.name == GOLDEN_FILE_NAME), None)
    test_files = [f for f in all_files if f.name != GOLDEN_FILE_NAME]

    if not golden_file:
        log(f"Golden file '{GOLDEN_FILE_NAME}' not found!", Colors.RED)
        sys.exit(1)

    # Run Tests
    passed, failed = 0, 0
    
    for test_file in test_files:
        log(f"\nComparing: {test_file.name} ...", Colors.BOLD)
        
        # Capture Success Boolean AND The Reason String
        is_success, reason_msg = compare_pptx_files(str(test_file), str(golden_file))
        
        if is_success:
            log(f"Result: PASS", Colors.GREEN)
            passed += 1
        else:
            # Print the specific reason in RED directly to console
            log(f"Result: FAIL", Colors.RED)
            log(f"  Reason: {reason_msg}", Colors.YELLOW)
            failed += 1

    # Summary
    log("\n" + "="*40)
    log(f"SUMMARY: Passed: {passed} | Failed: {failed}")
    log("="*40)
    log(f"Detailed logs saved to: {LOG_FILE}")
    log(f"Debug traces saved to: pptx_debug_{TIMESTAMP}.log")
    
    log_file_handle.close()
    sys.exit(1 if failed > 0 else 0)

if __name__ == "__main__":
    main()